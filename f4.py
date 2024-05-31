import os
from openpyxl import load_workbook
from docx import Document
from pptx import Presentation
from PIL import Image, ImageChops
import PyPDF2
from collections import defaultdict
from prettytable import PrettyTable
myTable1=PrettyTable(['File','Size(MB)'])
def get_excel_content(file_path):
    try:
        wb = load_workbook(file_path)
        all_content = []  # Initialize a list to store content from all sheets

        for sheet_name in wb.sheetnames:  # Iterate over all sheet names in the workbook
            sheet = wb[sheet_name]
            content = []

            for row in sheet.iter_rows(values_only=True):
                content.append(row)

            all_content.append(content)  # Append content of current sheet to all_content list

        return all_content  # Return the content from all sheets as a list of lists
    except Exception as e:
        print(f"Error reading file '{file_path}': {e}")
        return None


def find_duplicate_excel_files(directory):
    excel_files = {}
    duplicates = []

    for root, dirs, files in os.walk(directory):
        for file in files:
            if file.lower().endswith('.xlsx') or file.lower().endswith('.xlsb'):
                file_path = os.path.join(root, file)
                content = get_excel_content(file_path)
                if content:
                    content_tuple = tuple(map(tuple, content))  # Convert content list to tuple of tuples
                    if content_tuple in excel_files:
                        excel_files[content_tuple].append(file_path)
                    else:
                        excel_files[content_tuple] = [file_path]

    for file_paths in excel_files.values():
        if len(file_paths) > 1:
            duplicates.append(file_paths)

    return duplicates

def read_word_document(file_path):
    try:
        document = Document(file_path)
        text_content = "\n".join([paragraph.text for paragraph in document.paragraphs])
        return text_content.lower()
    except Exception as e:
        print(f"Error reading file '{file_path}': {e}")
        return None

def find_duplicate_word_files(directory):
    word_files = [file for file in os.listdir(directory) if file.lower().endswith('.docx')]
    duplicate_files = []

    for i, file1 in enumerate(word_files):
        for file2 in word_files[i+1:]:
            file1_path = os.path.join(directory, file1)
            file2_path = os.path.join(directory, file2)
            text1 = read_word_document(file1_path)
            text2 = read_word_document(file2_path)
            if text1 == text2:
                duplicate_files.append((file1_path,file2_path))

    return duplicate_files

def extract_first_three_words(slide):
    text = ''
    for shape in slide.shapes:
        if shape.has_text_frame:
            text += shape.text_frame.text.strip() + ' '
    words = text.split()
    first_three_words = ' '.join(words[:3])
    return first_three_words.lower()

'''def compare_slide_content(slide1, slide2):
    text1 = []
    text2 = []
    for slide 
    return text1 == text2'''

def compare_slide_content(slide1, slide2):
    # Compare the number of shapes in the slides
    if len(slide1.shapes) != len(slide2.shapes):
        return False

    # Compare the text content of shapes in the slides
    for shape1, shape2 in zip(slide1.shapes, slide2.shapes):
        if hasattr(shape1, "text") and hasattr(shape2, "text"):
            if shape1.text != shape2.text:
                return False

    return True

def compare_ppt(ppt_file1, ppt_file2):
    ppt1 = Presentation(ppt_file1)
    ppt2 = Presentation(ppt_file2)

    slides1 = ppt1.slides
    slides2 = ppt2.slides

    for slide1, slide2 in zip(slides1, slides2):
        if not compare_slide_content(slide1, slide2):
            return False

    return True

def are_images_identical(image_path1, image_path2):
    try:
        image1 = Image.open(image_path1)
        image2 = Image.open(image_path2)
        if image1.size != image2.size:
            return False
        diff = ImageChops.difference(image1, image2)
        return diff.getbbox() is None
    except IOError:
        print("Unable to open one or both images.")
        return False
def extract_pdf_content(file_path):
    try:
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            content = ''
            for page_num in range (len(reader.pages)):
                page = reader.pages[page_num]
                content += page.extract_text()
            return content.lower()
    except Exception as e:
        print(f"Error reading file '{file_path}': {e}")
        return None

def find_duplicate_pdf_files(directory):
    pdf_files = defaultdict(list)
    duplicates = []

    for root, _, files in os.walk(directory):
        for file in files:
            if file.lower().endswith('.pdf'):
                file_path = os.path.join(root, file)
                content = extract_pdf_content(file_path)
                if content:
                    pdf_files[content].append(file_path)

    for file_paths in pdf_files.values():
        if len(file_paths) > 1:
            duplicates.append(file_paths)

    return duplicates

# Directory to scan for duplicate files
directory = input("Enter the directory path: ")

# Find and print duplicate Excel files
duplicate_excel_files = find_duplicate_excel_files(directory)
if duplicate_excel_files:
    print("Duplicate Excel files found:")
    for files in duplicate_excel_files:
        print(files)
        for file_path in files:
            
            s = os.path.getsize(file_path) / (1024 * 1024)  # Get the size of each individual file
            
            myTable1.add_row([file_path, s])
            
         
#else:
    #print("No duplicate Excel files found.")

# Find and print duplicate Word files
duplicate_word_files = find_duplicate_word_files(directory)
if duplicate_word_files:
    print("Duplicate Word documents found:")
    for file1, file2 in duplicate_word_files:
        print(f"Duplicate pair: {file1} and {file2}")
        s=os.path.getsize(file1)/(1024*1024)
        s2 = os.path.getsize(file2) / (1024 * 1024)
        myTable1.add_row([file1,s])
        myTable1.add_row([file2,s2])
#else:
    #print("No duplicate Word documents found.")

# Find and print duplicate PPT files
ppt_files = [file for file in os.listdir(directory) if file.lower().endswith('.pptx')]
for i, file1 in enumerate(ppt_files):
    for file2 in ppt_files[i+1:]:
        ppt_file1 = os.path.join(directory, file1)
        ppt_file2 = os.path.join(directory, file2)
        if compare_ppt(ppt_file1, ppt_file2):
            print(f"The PPT files '{file1}' and '{file2}' have the same text content on slides.")
            s=os.path.getsize(ppt_file1)/(1024*1024)
            myTable1.add_row([file1,s])
        #else:
            #print(f"The PPT files '{file1}' and '{file2}' have different text content on slides.")

# Find and print duplicate images
image_files = [file for file in os.listdir(directory) if file.lower().endswith(('.jpg', '.jpeg', '.png', '.bmp'))]
for i, file1 in enumerate(image_files):
    for file2 in image_files[i+1:]:
        image1_path = os.path.join(directory, file1)
        image2_path = os.path.join(directory, file2)
        if are_images_identical(image1_path, image2_path):
            print(f"The images '{file1}' and '{file2}' are identical.")
            s=os.path.getsize(image1_path)/(1024*1024)
            myTable1.add_row([file1,s])
        #else:
           # print(f"The images '{file1}' and '{file2}' are different.")
pdf_files = [file for file in os.listdir(directory) if file.lower().endswith('.pdf')]
for i, file1 in enumerate(pdf_files):
    for file2 in pdf_files[i+1:]:
        pdf1_path = os.path.join(directory, file1)
        pdf2_path = os.path.join(directory, file2)
        if find_duplicate_pdf_files(directory):
            print(f"The pdfs '{file1}' and '{file2}' are identical.")
            for file_path in files:
             s = os.path.getsize(file_path) / (1024 * 1024)
             myTable1.add_row([file_path, s])
        #else:
           # print(f"The images '{file1}' and '{file2}' are different.")
print(myTable1)
a=input()
'''if duplicate_word_files:
    option=(input("Do you want to delete word files,  yes or no"))
    option=option.lower()
    if option=='yes':
        for file in duplicate_word_files:
            s.'''